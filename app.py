import streamlit as st
import os
import tempfile
from pathlib import Path
import io
from PIL import Image
import subprocess
import shutil

# PowerPoint to PDF imports
try:
    from pptx import Presentation
    from reportlab.lib.pagesizes import A4
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# PDF to TIFF imports
try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

# Page configuration
st.set_page_config(
    page_title="Ijazul Haq | File Converter",
    page_icon="🔄",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better styling
st.markdown("""
    <style>
    .main-header {
        text-align: center;
        padding: 2rem 0;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 10px;
        margin-bottom: 2rem;
    }
    .stButton>button {
        width: 100%;
        background-color: #667eea;
        color: white;
        font-weight: bold;
        border-radius: 5px;
        padding: 0.5rem 1rem;
    }
    .stButton>button:hover {
        background-color: #764ba2;
    }
    .info-box {
        padding: 1rem;
        background-color: #f0f2f6;
        border-radius: 5px;
        margin: 1rem 0;
    }
    .footer {
        text-align: center;
        padding: 2rem 0;
        color: #666;
        margin-top: 3rem;
    }
    </style>
""", unsafe_allow_html=True)

# Header
st.markdown("""
    <div class="main-header">
        <h1>Ijazul Haq | File Converter</h1>
        <p>Comparative & Evolutionary Genomics Lab</p>
        <p style="font-size: 0.85em; margin-top: 0.3rem;">Dr. Amir Ali Abbasi | Principal Investigator of <a href="https://www.pahgncb.com/genomedb/public/" target="_blank" style="color: #fff; text-decoration: underline;">PAHG database</a></p>
        <p style="font-size: 0.9em; margin-top: 0.5rem;">Convert PowerPoint to PDF and PDF to TIFF with ease</p>
    </div>
""", unsafe_allow_html=True)

# Initialize session state
if 'last_dpi' not in st.session_state:
    st.session_state.last_dpi = 300
if 'last_compression' not in st.session_state:
    st.session_state.last_compression = 'tiff_deflate'

def convert_pptx_to_pdf_libreoffice(uploaded_file):
    """Convert PowerPoint file to PDF using LibreOffice (better visual rendering)"""
    try:
        import subprocess
        import shutil
        
        # Check if LibreOffice is available
        libreoffice_path = shutil.which('libreoffice') or shutil.which('soffice')
        if not libreoffice_path:
            raise ValueError("LibreOffice is not installed. Install it with: brew install --cask libreoffice")
        
        # Reset file pointer
        uploaded_file.seek(0)
        
        # Create temporary directory for conversion
        with tempfile.TemporaryDirectory() as tmp_dir:
            # Save uploaded file
            input_path = os.path.join(tmp_dir, uploaded_file.name)
            with open(input_path, 'wb') as f:
                f.write(uploaded_file.read())
            
            # Convert using LibreOffice headless mode
            output_dir = tmp_dir
            cmd = [
                libreoffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', output_dir,
                input_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode != 0:
                raise ValueError(f"LibreOffice conversion failed: {result.stderr}")
            
            # Find the generated PDF
            pdf_path = os.path.join(output_dir, Path(uploaded_file.name).stem + '.pdf')
            if not os.path.exists(pdf_path):
                raise ValueError("PDF file was not generated")
            
            # Read PDF
            with open(pdf_path, 'rb') as f:
                pdf_bytes = f.read()
            
            return pdf_bytes
    
    except Exception as e:
        st.error(f"Error converting with LibreOffice: {str(e)}")
        return None

def convert_pptx_to_pdf_api(uploaded_file, api_key=None):
    """Convert PowerPoint file to PDF using external API service"""
    try:
        import requests
        
        if not api_key:
            # Try using a free service or show instructions
            st.warning("⚠️ API key not provided. Using LibreOffice method instead.")
            return convert_pptx_to_pdf_libreoffice(uploaded_file)
        
        # Example using CloudConvert API (you can replace with any API service)
        uploaded_file.seek(0)
        files = {'file': (uploaded_file.name, uploaded_file.read(), 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
        
        # This is a placeholder - replace with actual API endpoint
        # For example, CloudConvert: https://api.cloudconvert.com/v2/convert
        response = requests.post(
            'https://api.cloudconvert.com/v2/convert',
            headers={'Authorization': f'Bearer {api_key}'},
            files=files,
            data={'format': 'pdf'}
        )
        
        if response.status_code != 200:
            raise ValueError(f"API conversion failed: {response.text}")
        
        return response.content
    
    except Exception as e:
        st.error(f"Error converting with API: {str(e)}")
        return None

def convert_pptx_to_pdf(uploaded_file, method='text'):
    """Convert PowerPoint file to PDF using specified method"""
    if method == 'api':
        # Check for API key in session state or environment
        api_key = st.session_state.get('api_key', None) or os.getenv('CONVERSION_API_KEY')
        return convert_pptx_to_pdf_api(uploaded_file, api_key)
    elif method == 'libreoffice':
        return convert_pptx_to_pdf_libreoffice(uploaded_file)
    else:
        # Default text extraction method
        return convert_pptx_to_pdf_text(uploaded_file)

def convert_pptx_to_pdf_text(uploaded_file):
    """Convert PowerPoint file to PDF using text extraction (original method)"""
    try:
        # Reset file pointer
        uploaded_file.seek(0)
        
        # Check file extension
        file_ext = Path(uploaded_file.name).suffix.lower()
        if file_ext == '.ppt':
            st.warning("⚠️ Note: .ppt files (older format) may not work perfectly. For best results, use .pptx files.")
        
        # Save uploaded file to temporary location
        suffix = '.pptx' if file_ext in ['.ppt', '.pptx'] else '.pptx'
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_file:
            tmp_file.write(uploaded_file.read())
            tmp_path = tmp_file.name
        
        # Load presentation (python-pptx only supports .pptx)
        if file_ext == '.ppt':
            raise ValueError("The .ppt format (PowerPoint 97-2003) is not directly supported. Please convert your file to .pptx format first, or use LibreOffice to convert it.")
        
        prs = Presentation(tmp_path)
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        
        doc = SimpleDocTemplate(pdf_buffer, pagesize=A4,
                               rightMargin=72, leftMargin=72,
                               topMargin=72, bottomMargin=18)
        
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=16,
            textColor='#333333',
            spaceAfter=30,
            alignment=TA_CENTER
        )
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=11,
            textColor='#000000',
            spaceAfter=12,
            alignment=TA_LEFT
        )
        
        story = []
        
        # Convert each slide
        for i, slide in enumerate(prs.slides):
            if i > 0:
                story.append(PageBreak())
            
            # Add slide title
            story.append(Paragraph(f"Slide {i + 1}", title_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Extract and add text content
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            if text_content:
                for text in text_content:
                    # Clean and format text
                    clean_text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    # Split into paragraphs
                    paragraphs = clean_text.split('\n')
                    for para in paragraphs:
                        if para.strip():
                            story.append(Paragraph(para.strip(), body_style))
                            story.append(Spacer(1, 0.1*inch))
            else:
                story.append(Paragraph("(No text content on this slide)", body_style))
        
        # Build PDF
        doc.build(story)
        pdf_buffer.seek(0)
        
        # Clean up
        os.unlink(tmp_path)
        
        return pdf_buffer.getvalue()
    
    except Exception as e:
        st.error(f"Error converting PowerPoint to PDF: {str(e)}")
        import traceback
        st.error(f"Details: {traceback.format_exc()}")
        return None

def check_poppler_available():
    """Check if Poppler is available and return path if found"""
    try:
        import subprocess
        # Check environment variable first
        poppler_path_env = os.getenv('POPPLER_PATH')
        if poppler_path_env:
            pdftoppm_path = os.path.join(poppler_path_env, 'pdftoppm')
            if os.path.exists(pdftoppm_path):
                # Test if it works
                try:
                    result = subprocess.run([pdftoppm_path, '-v'], capture_output=True, timeout=5)
                    output = (result.stdout.decode() + result.stderr.decode()).lower()
                    if 'version' in output or result.returncode == 0:
                        return pdftoppm_path
                except:
                    pass
        
        # Try common Poppler paths
        poppler_paths = [
            shutil.which('pdftoppm'),
            shutil.which('pdftocairo'),
            '/usr/bin/pdftoppm',
            '/usr/local/bin/pdftoppm',
            '/opt/homebrew/bin/pdftoppm',  # macOS Apple Silicon
            '/app/.apt/usr/bin/pdftoppm',  # Streamlit Cloud
            '/app/.apt/usr/bin/pdftocairo',  # Streamlit Cloud alternative
            '/usr/bin/pdftocairo',  # Alternative binary
        ]
        
        for path in poppler_paths:
            if path and os.path.exists(path):
                # Test if it works - try multiple ways
                try:
                    # Try with -v flag
                    result = subprocess.run([path, '-v'], capture_output=True, timeout=5)
                    output = (result.stdout.decode() + result.stderr.decode()).lower()
                    if 'version' in output or result.returncode == 0:
                        return path
                except Exception as e1:
                    try:
                        # Try with --version flag
                        result = subprocess.run([path, '--version'], capture_output=True, timeout=5)
                        output = (result.stdout.decode() + result.stderr.decode()).lower()
                        if 'version' in output or result.returncode == 0:
                            return path
                    except:
                        # If both fail but file exists, still return it (might work for conversion)
                        # Only do this if it's from shutil.which (in PATH)
                        if path == shutil.which('pdftoppm') or path == shutil.which('pdftocairo'):
                            return path
                        continue
        
        return None
    except Exception as e:
        # If there's an error, try a simple which check as fallback
        try:
            simple_path = shutil.which('pdftoppm')
            if simple_path and os.path.exists(simple_path):
                return simple_path
        except:
            pass
        return None

def convert_pdf_to_tiff(pdf_bytes, dpi=300, compression='tiff_deflate'):
    """Convert PDF to multi-page TIFF"""
    try:
        # Check Poppler availability first (but don't fail if detection fails on cloud)
        poppler_path = check_poppler_available()
        
        # On Streamlit Cloud, even if detection fails, try conversion anyway
        # since Poppler might be installed but subprocess calls are restricted
        if not poppler_path:
            # Check if we're likely on Streamlit Cloud
            is_streamlit_cloud = os.path.exists('/app/.apt/usr/bin/pdftoppm') or os.path.exists('/app/.apt/usr/bin/pdftocairo')
            
            if not is_streamlit_cloud:
                error_msg = """
                **Poppler is not installed or not in PATH.**
                
                **For local development:**
                - macOS: `brew install poppler`
                - Linux: `sudo apt-get install poppler-utils`
                - Windows: Download from [poppler-windows](https://github.com/oschwartz10612/poppler-windows/releases)
                
                **For cloud deployments (Streamlit Cloud, Heroku, etc.):**
                Poppler needs to be installed in your deployment environment. Options:
                1. Add Poppler to your buildpack/dockerfile
                2. Use a Docker image with Poppler pre-installed
                3. Set POPPLER_PATH environment variable if Poppler is in a custom location
                """
                st.error(error_msg)
                return None
            else:
                # On Streamlit Cloud, try anyway - Poppler might be installed but detection failed
                st.info("⚠️ Poppler detection failed, but attempting conversion anyway (Poppler may still be available)...")
                poppler_path = None  # Will use default path
        
        # Convert PDF to images
        # Try to use poppler_path if available
        poppler_dir = None
        if poppler_path:
            poppler_dir = os.path.dirname(poppler_path)
            # If poppler_path is a directory, use it directly
            if os.path.isdir(poppler_path):
                poppler_dir = poppler_path
        
        # For Streamlit Cloud, try common paths even if detection failed
        if not poppler_dir:
            streamlit_cloud_paths = ['/app/.apt/usr/bin', '/usr/bin']
            for potential_dir in streamlit_cloud_paths:
                if os.path.exists(os.path.join(potential_dir, 'pdftoppm')):
                    poppler_dir = potential_dir
                    break
        
        try:
            if poppler_dir:
                images = convert_from_bytes(pdf_bytes, dpi=dpi, poppler_path=poppler_dir)
            else:
                images = convert_from_bytes(pdf_bytes, dpi=dpi)
        except Exception as e:
            # Fallback: try Streamlit Cloud path explicitly
            if not poppler_dir and os.path.exists('/app/.apt/usr/bin/pdftoppm'):
                try:
                    images = convert_from_bytes(pdf_bytes, dpi=dpi, poppler_path='/app/.apt/usr/bin')
                except:
                    raise e
            else:
                raise e
        
        if not images:
            return None
        
        # Save as multi-page TIFF
        tiff_buffer = io.BytesIO()
        
        # Determine compression
        if compression == 'tiff_deflate':
            comp = 'tiff_deflate'
        elif compression == 'tiff_lzw':
            comp = 'tiff_lzw'
        elif compression == 'tiff_adobe_deflate':
            comp = 'tiff_adobe_deflate'
        else:
            comp = None
        
        # Save first image to get format
        if len(images) == 1:
            images[0].save(tiff_buffer, format='TIFF', compression=comp)
        else:
            images[0].save(
                tiff_buffer,
                format='TIFF',
                compression=comp,
                save_all=True,
                append_images=images[1:]
            )
        
        tiff_buffer.seek(0)
        return tiff_buffer.getvalue()
    
    except Exception as e:
        error_str = str(e)
        if "poppler" in error_str.lower() or "page count" in error_str.lower():
            error_msg = f"""
            **Poppler Error: {error_str}**
            
            Poppler is required for PDF to TIFF conversion but is not properly configured.
            
            **Troubleshooting:**
            1. Check if Poppler is installed: `pdftoppm -v` (should show version)
            2. For cloud deployments, ensure Poppler is included in your deployment configuration
            3. Set POPPLER_PATH environment variable if Poppler is in a non-standard location
            4. Check that Poppler binaries are executable
            
            **For Streamlit Cloud:**
            - Add Poppler installation to your `packages.txt` or use a Dockerfile
            - Example packages.txt: `poppler-utils`
            """
            st.error(error_msg)
        else:
            st.error(f"Error converting PDF to TIFF: {error_str}")
        return None

def get_pdf_preview(pdf_bytes):
    """Get first page of PDF as image for preview"""
    try:
        poppler_path = check_poppler_available()
        poppler_dir = None
        if poppler_path:
            poppler_dir = os.path.dirname(poppler_path)
            if os.path.isdir(poppler_path):
                poppler_dir = poppler_path
        
        try:
            if poppler_dir:
                images = convert_from_bytes(
                    pdf_bytes, 
                    dpi=150, 
                    first_page=1, 
                    last_page=1,
                    poppler_path=poppler_dir
                )
            else:
                images = convert_from_bytes(pdf_bytes, dpi=150, first_page=1, last_page=1)
        except:
            images = convert_from_bytes(pdf_bytes, dpi=150, first_page=1, last_page=1)
        
        if images:
            return images[0]
        return None
    except Exception as e:
        # Don't show error for preview failures, just return None
        return None

def get_tiff_preview(tiff_bytes):
    """Get first page of TIFF as image for preview"""
    try:
        img = Image.open(io.BytesIO(tiff_bytes))
        return img
    except Exception as e:
        st.warning(f"Could not generate TIFF preview: {str(e)}")
        return None

# Sidebar for navigation
st.sidebar.title("Navigation")
tool_choice = st.sidebar.radio(
    "Select Tool",
    ["PowerPoint to PDF", "PDF to TIFF"],
    index=0
)

# Main content area
if tool_choice == "PowerPoint to PDF":
    st.header("📄 PowerPoint to PDF Converter")
    
    if not PPTX_AVAILABLE:
        st.error("⚠️ Required libraries not installed. Please install python-pptx and reportlab.")
        st.code("pip install python-pptx reportlab", language="bash")
    else:
        uploaded_file = st.file_uploader(
            "Choose a PowerPoint file",
            type=['ppt', 'pptx'],
            help="Upload a .ppt or .pptx file to convert to PDF"
        )
        
        if uploaded_file is not None:
            # Show file info
            uploaded_file.seek(0)
            file_size = len(uploaded_file.read())
            uploaded_file.seek(0)
            st.info(f"📁 File: {uploaded_file.name} ({file_size / 1024:.2f} KB)")
            
            # Conversion method selector
            conversion_methods = ['Text Extraction', 'LibreOffice (API)', 'External API']
            method = st.selectbox(
                "Conversion Method",
                conversion_methods,
                help="Text Extraction: Fast, extracts text only. LibreOffice: Better visual rendering (requires LibreOffice). External API: Uses cloud service (requires API key)."
            )
            
            # API key input if External API is selected
            api_key = None
            if method == 'External API':
                api_key = st.text_input(
                    "API Key (optional)",
                    type="password",
                    help="Enter your conversion API key. Leave empty to use LibreOffice method instead.",
                    value=st.session_state.get('api_key', '')
                )
                if api_key:
                    st.session_state['api_key'] = api_key
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("🔄 Convert to PDF", type="primary"):
                    with st.spinner("Converting PowerPoint to PDF... This may take a moment."):
                        uploaded_file.seek(0)
                        
                        # Map method name to function parameter
                        method_param = 'text'
                        if method == 'LibreOffice (API)':
                            method_param = 'libreoffice'
                        elif method == 'External API':
                            method_param = 'api'
                        
                        pdf_bytes = convert_pptx_to_pdf(uploaded_file, method=method_param)
                        
                        if pdf_bytes:
                            st.session_state['pdf_output'] = pdf_bytes
                            st.session_state['pdf_filename'] = Path(uploaded_file.name).stem + ".pdf"
                            output_size = len(pdf_bytes)
                            st.success(f"✅ Conversion successful! Output size: {output_size / 1024:.2f} KB")
            
            if 'pdf_output' in st.session_state:
                with col2:
                    st.download_button(
                        label="📥 Download PDF",
                        data=st.session_state['pdf_output'],
                        file_name=st.session_state['pdf_filename'],
                        mime="application/pdf",
                        type="primary"
                    )
                
                # Preview
                st.subheader("📊 Preview")
                preview_img = get_pdf_preview(st.session_state['pdf_output'])
                if preview_img:
                    st.image(preview_img, caption="First page preview", use_container_width=True)
                else:
                    st.info("Preview not available for this PDF")

elif tool_choice == "PDF to TIFF":
    st.header("🖼️ PDF to TIFF Converter")
    
    if not PDF2IMAGE_AVAILABLE:
        st.error("⚠️ Required libraries not installed. Please install pdf2image and Pillow.")
        st.code("pip install pdf2image pillow", language="bash")
        st.info("""
        **Note:** pdf2image requires Poppler. Install it:
        - **macOS:** `brew install poppler`
        - **Linux:** `sudo apt-get install poppler-utils`
        - **Windows:** Download from [poppler-windows](https://github.com/oschwartz10612/poppler-windows/releases)
        """)
    else:
        # Check Poppler status
        poppler_status = check_poppler_available()
        if poppler_status:
            st.success(f"✅ Poppler is available")
        else:
            # On Streamlit Cloud, try to be more lenient - the conversion might still work
            # even if detection fails due to subprocess restrictions
            st.info("ℹ️ Poppler detection: Checking availability...")
            st.info("""
            **Note:** If you're on Streamlit Cloud and have `packages.txt` with `poppler-utils`, 
            Poppler should be installed. The conversion will be attempted even if detection shows a warning.
            """)
            with st.expander("📋 Troubleshooting Poppler on Streamlit Cloud"):
                st.markdown("""
                **For Streamlit Cloud:**
                1. Ensure `packages.txt` exists in your repo root with:
                   ```
                   poppler-utils
                   ```
                2. Redeploy your app after adding/updating packages.txt
                3. Poppler should be installed at `/app/.apt/usr/bin/pdftoppm`
                
                **If conversion still fails:**
                - Check Streamlit Cloud logs for Poppler installation errors
                - Verify packages.txt is in the root directory (same level as app.py)
                - Try redeploying the app
                """)
        
        uploaded_file = st.file_uploader(
            "Choose a PDF file",
            type=['pdf'],
            help="Upload a .pdf file to convert to multi-page TIFF"
        )
        
        if uploaded_file is not None:
            # Show file info
            uploaded_file.seek(0)
            file_size = len(uploaded_file.read())
            uploaded_file.seek(0)
            st.info(f"📁 File: {uploaded_file.name} ({file_size / 1024:.2f} KB)")
            
            # Options
            col1, col2 = st.columns(2)
            
            with col1:
                dpi_options = [150, 300, 600, 700, 800, 900, 1000, 1100, 1200]
                dpi = st.selectbox(
                    "DPI (Resolution)",
                    dpi_options,
                    index=dpi_options.index(st.session_state.last_dpi) if st.session_state.last_dpi in dpi_options else 1,
                    help="Higher DPI = better quality but larger file size. 300 DPI is recommended for most uses."
                )
                st.session_state.last_dpi = dpi
            
            with col2:
                compression = st.selectbox(
                    "Compression Method",
                    ['tiff_deflate', 'tiff_lzw', 'tiff_adobe_deflate', 'None'],
                    index=['tiff_deflate', 'tiff_lzw', 'tiff_adobe_deflate', 'None'].index(st.session_state.last_compression) if st.session_state.last_compression in ['tiff_deflate', 'tiff_lzw', 'tiff_adobe_deflate', 'None'] else 0,
                    help="Compression reduces file size. Deflate is recommended for best balance."
                )
                st.session_state.last_compression = compression
            
            # Convert button
            if st.button("🔄 Convert to TIFF", type="primary"):
                with st.spinner("Converting PDF to TIFF... This may take a moment."):
                    pdf_bytes = uploaded_file.read()
                    tiff_bytes = convert_pdf_to_tiff(
                        pdf_bytes,
                        dpi=dpi,
                        compression=compression if compression != 'None' else None
                    )
                    
                    if tiff_bytes:
                        st.session_state['tiff_output'] = tiff_bytes
                        st.session_state['tiff_filename'] = Path(uploaded_file.name).stem + ".tiff"
                        st.success("✅ Conversion successful!")
            
            # Download button and preview
            if 'tiff_output' in st.session_state:
                output_size = len(st.session_state['tiff_output'])
                st.info(f"📁 Output file size: {output_size / 1024:.2f} KB")
                
                st.download_button(
                    label="📥 Download TIFF",
                    data=st.session_state['tiff_output'],
                    file_name=st.session_state['tiff_filename'],
                    mime="image/tiff",
                    type="primary"
                )
                
                # Preview
                st.subheader("📊 Preview")
                preview_img = get_tiff_preview(st.session_state['tiff_output'])
                if preview_img:
                    st.image(preview_img, caption="First page preview", use_container_width=True)
                else:
                    st.info("Preview not available for this TIFF")

# Footer
st.markdown("""
    <div class="footer">
        <hr>
        <p><strong>Ijazul Haq | File Converter</strong></p>
        <p>Comparative & Evolutionary Genomics Lab</p>
        <p style="font-size: 0.8em; margin-top: 0.5rem;">© 2024 | Built with Streamlit</p>
        <p style="font-size: 0.8em;">Supports: PowerPoint (PPT/PPTX) → PDF | PDF → TIFF</p>
    </div>
""", unsafe_allow_html=True)

